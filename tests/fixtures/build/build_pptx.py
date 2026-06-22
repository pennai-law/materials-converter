"""Build tests/fixtures/sample_with_notes.pptx — a 3-slide deck where
slides 1 and 2 have speaker notes and slide 3 does not. Used by:
  - The default PPTX conversion test (slide markers + speaker-notes markers)
  - The --notes-only test (transcript skips note-less slides)

Run from repo root:
    ./venv/bin/python tests/fixtures/build/build_pptx.py
"""
from pathlib import Path

from pptx import Presentation

OUT = Path(__file__).resolve().parent.parent / "sample_with_notes.pptx"


def build() -> Path:
    prs = Presentation()

    s1 = prs.slides.add_slide(prs.slide_layouts[1])  # Title and Content
    s1.shapes.title.text = "Patent Infringement Doctrine"
    s1.placeholders[1].text = (
        "- Literal infringement\n"
        "- Doctrine of equivalents\n"
        "- All-elements rule"
    )
    s1.notes_slide.notes_text_frame.text = (
        "Speaker note for slide 1: walk through the Warner-Jenkinson "
        "framework. Emphasize that the doctrine of equivalents extends "
        "literal infringement when the elements are equivalent."
    )

    s2 = prs.slides.add_slide(prs.slide_layouts[1])
    s2.shapes.title.text = "Festo and Prosecution History Estoppel"
    s2.placeholders[1].text = (
        "- Festo Corp. v. Shoketsu (2002)\n"
        "- Narrowing amendments\n"
        "- Rebuttable presumption of estoppel"
    )
    s2.notes_slide.notes_text_frame.text = (
        "Speaker note for slide 2: Festo cabined the doctrine of "
        "equivalents — narrowing amendments raise a rebuttable presumption "
        "that the patentee surrendered the equivalent territory."
    )

    s3 = prs.slides.add_slide(prs.slide_layouts[1])
    s3.shapes.title.text = "Summary and Discussion"
    s3.placeholders[1].text = (
        "- Three takeaways\n"
        "- Open questions for next class"
    )
    # Intentionally NO speaker notes on slide 3 — tests verify the marker
    # is absent for note-less slides.

    prs.save(str(OUT))
    return OUT


if __name__ == "__main__":
    p = build()
    size_kb = p.stat().st_size / 1024
    print(f"Wrote {p} ({size_kb:.1f} KB)")
